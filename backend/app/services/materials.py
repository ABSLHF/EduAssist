from __future__ import annotations

import io
import logging
from datetime import datetime
from pathlib import Path

from fastapi import UploadFile

from app.config import settings

try:
    import oss2  # type: ignore
except Exception:
    oss2 = None

# Optional parsers
try:
    import pdfplumber  # type: ignore
except Exception:
    pdfplumber = None

try:
    from docx import Document  # type: ignore
except Exception:
    Document = None

try:
    from pptx import Presentation  # type: ignore
except Exception:
    Presentation = None

logger = logging.getLogger(__name__)


def _safe_filename(filename: str) -> str:
    raw = Path(filename or "upload.bin").name
    cleaned = "".join(ch if ch.isalnum() or ch in "._-()[] " else "_" for ch in raw).strip()
    return cleaned or "upload.bin"


def _oss_endpoint_url() -> str:
    endpoint = (settings.oss_endpoint or "").strip()
    if not endpoint:
        raise RuntimeError("OSS endpoint is not configured")
    if endpoint.startswith("http://") or endpoint.startswith("https://"):
        return endpoint
    scheme = "https" if settings.oss_secure else "http"
    return f"{scheme}://{endpoint}"


def _create_oss_bucket(bucket_name: str | None = None):
    if oss2 is None:
        raise RuntimeError("oss2 package is not installed")
    access_key_id = (settings.oss_access_key_id or "").strip()
    access_key_secret = (settings.oss_access_key_secret or "").strip()
    bucket = (bucket_name or settings.oss_bucket or "").strip()
    if not access_key_id or not access_key_secret or not bucket:
        raise RuntimeError("OSS credentials or bucket are not configured")
    auth = oss2.Auth(access_key_id, access_key_secret)
    return oss2.Bucket(auth, _oss_endpoint_url(), bucket)


def _build_object_key(course_id: int, filename: str) -> str:
    stamp = datetime.utcnow().strftime("%Y%m%d%H%M%S%f")
    return f"course_{course_id}/{stamp}_{filename}"


def _save_local(course_id: int, filename: str, data: bytes, base_dir: str = "storage") -> str:
    course_dir = Path(base_dir) / f"course_{course_id}"
    course_dir.mkdir(parents=True, exist_ok=True)
    dest = course_dir / filename
    with dest.open("wb") as f:
        f.write(data)
    return str(dest)


def _save_oss(course_id: int, filename: str, data: bytes) -> str:
    bucket = _create_oss_bucket()
    object_key = _build_object_key(course_id, filename)
    bucket.put_object(object_key, data)
    return f"oss://{bucket.bucket_name}/{object_key}"


def save_upload(course_id: int, upload: UploadFile, base_dir: str = "storage") -> str:
    filename = _safe_filename(upload.filename or "upload.bin")
    data = upload.file.read()
    if not data:
        raise RuntimeError("Uploaded file is empty")
    return save_upload_bytes(course_id=course_id, filename=filename, data=data, base_dir=base_dir)


def save_upload_bytes(course_id: int, filename: str, data: bytes, base_dir: str = "storage") -> str:
    mode = (settings.material_storage_mode or "local").strip().lower().replace("-", "_")
    prefer_oss = mode in {"oss", "oss_only", "oss_first", "oss_fallback", "oss_fallback_local"}
    allow_local = mode in {"", "local", "oss_first", "oss_fallback", "oss_fallback_local"}

    safe_name = _safe_filename(filename)
    if prefer_oss:
        try:
            return _save_oss(course_id=course_id, filename=safe_name, data=data)
        except Exception as exc:
            logger.warning("OSS upload failed, fallback to local storage: %s", exc)
            if not allow_local:
                raise

    if allow_local:
        return _save_local(course_id=course_id, filename=safe_name, data=data, base_dir=base_dir)

    raise RuntimeError(f"Unsupported MATERIAL_STORAGE_MODE: {settings.material_storage_mode}")


def _parse_oss_ref(file_ref: str) -> tuple[str, str] | None:
    marker = "oss://"
    if not file_ref.startswith(marker):
        return None
    raw = file_ref[len(marker) :]
    bucket, sep, key = raw.partition("/")
    if not sep or not bucket or not key:
        raise RuntimeError(f"Invalid OSS file reference: {file_ref}")
    return bucket, key


def _resolve_local_path(file_ref: str) -> Path:
    path = Path(file_ref)
    if path.is_absolute():
        return path
    backend_root = Path(__file__).resolve().parents[2]
    candidates = [path, backend_root / file_ref]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return candidates[0]


def resolve_material_filename(file_ref: str, fallback_title: str | None = None) -> str:
    parsed = _parse_oss_ref(file_ref)
    if parsed is not None:
        _, key = parsed
        name = Path(key).name
        if name:
            return name
    name = Path(file_ref).name
    if name:
        return name
    return _safe_filename(fallback_title or "downloaded_file")


def read_file_bytes(file_ref: str) -> bytes:
    parsed = _parse_oss_ref(file_ref)
    if parsed is not None:
        bucket_name, key = parsed
        bucket = _create_oss_bucket(bucket_name=bucket_name)
        obj = bucket.get_object(key)
        return obj.read()

    path = _resolve_local_path(file_ref)
    if not path.exists():
        raise FileNotFoundError(f"Material file not found: {file_ref}")
    return path.read_bytes()


def extract_text_from_bytes(data: bytes, suffix: str) -> str:
    lower_suffix = suffix.lower()
    if lower_suffix in [".txt", ".md"]:
        return data.decode("utf-8", errors="ignore")
    if lower_suffix == ".pdf" and pdfplumber is not None:
        text_parts = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        return "\n".join(text_parts)
    if lower_suffix in [".docx", ".doc"] and Document is not None:
        doc = Document(io.BytesIO(data))
        return "\n".join([p.text for p in doc.paragraphs if p.text])
    if lower_suffix in [".pptx", ".ppt"] and Presentation is not None:
        prs = Presentation(io.BytesIO(data))
        slides_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slides_text.append(shape.text)
        return "\n".join(slides_text)
    return data.decode("utf-8", errors="ignore")


def extract_text(file_path: str) -> str:
    file_name = resolve_material_filename(file_path)
    suffix = Path(file_name).suffix.lower()
    data = read_file_bytes(file_path)
    return extract_text_from_bytes(data=data, suffix=suffix)
